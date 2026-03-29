"""
Standalone PPTX text extraction script.

Usage:
    python pptx_extract.py <pptx_file_url>

Downloads the PPTX file from the given URL, extracts all text from all slides
(including shapes, tables, and grouped shapes), and prints the text to stdout.
Errors are written to stderr and the process exits with code 1 on failure.
"""

import sys
import os
import io
import zipfile
import tempfile
import requests
from pptx import Presentation
from pptx.exc import PackageNotFoundError

# Force UTF-8 output on Windows (avoids 'charmap' codec errors with Unicode text)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")


def extract_text_from_shape(shape) -> list[str]:
    """Recursively extract text from a shape, handling groups and tables."""
    texts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                texts.append(para_text)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    texts.append(cell_text)

    # Recurse into grouped shapes
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            texts.extend(extract_text_from_shape(child_shape))

    return texts


def extract_slides(prs: Presentation) -> list[str]:
    """Extract text from all slides, prefixed with slide number headers."""
    all_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            slide_texts.extend(extract_text_from_shape(shape))

        if slide_texts:
            all_text.append(f"--- Slide {slide_num} ---")
            all_text.extend(slide_texts)

    return all_text


def extract_notes(prs: Presentation) -> list[str]:
    """Extract speaker notes from all slides that have them."""
    notes_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text.strip()
            if text:
                notes_text.append(f"--- Slide {slide_num} Notes ---")
                notes_text.append(text)

    return notes_text


def download_file(url: str) -> str:
    """Download a file to a temp path and return that path."""
    response = requests.get(url, timeout=120)
    response.raise_for_status()

    fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    with os.fdopen(fd, "wb") as f:
        f.write(response.content)

    return tmp_path


def main() -> None:
    if len(sys.argv) != 2:
        print("Usage: python pptx_extract.py <pptx_file_url>", file=sys.stderr)
        sys.exit(1)

    file_url = sys.argv[1]
    tmp_path = None

    try:
        tmp_path = download_file(file_url)
        prs = Presentation(tmp_path)

        lines = extract_slides(prs) + extract_notes(prs)
        print("\n".join(lines))

    except requests.exceptions.Timeout:
        print("Download timed out after 120 seconds.", file=sys.stderr)
        sys.exit(1)
    except requests.exceptions.HTTPError as e:
        print(f"Failed to download file (HTTP {e.response.status_code}): {e}", file=sys.stderr)
        sys.exit(1)
    except requests.RequestException as e:
        print(f"Failed to download file: {e}", file=sys.stderr)
        sys.exit(1)
    except zipfile.BadZipFile:
        print(
            "File is not a valid PPTX (bad ZIP structure). "
            "Legacy .ppt format is not supported — please convert to .pptx.",
            file=sys.stderr,
        )
        sys.exit(1)
    except PackageNotFoundError:
        print(
            "File is not a valid PPTX package. "
            "Ensure the file is a .pptx (PowerPoint Open XML) file.",
            file=sys.stderr,
        )
        sys.exit(1)
    except Exception as e:
        print(f"PPTX extraction failed: {e}", file=sys.stderr)
        sys.exit(1)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)


if __name__ == "__main__":
    main()
